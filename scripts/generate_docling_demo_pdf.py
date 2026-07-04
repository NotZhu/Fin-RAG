"""生成用于 Docling 解析演示的多格式企业金融文档。"""

from __future__ import annotations

import argparse
import csv
import json
import tempfile
import zipfile
from copy import copy
from html import escape
from pathlib import Path

import fitz
from PIL import Image, ImageDraw, ImageFont


PROJECT_ROOT = Path(__file__).resolve().parents[1]
DEFAULT_OUTPUT_DIR = PROJECT_ROOT / "output" / "demo-documents"
DEFAULT_BASENAME = "finrag_docling_demo"
FONT_CANDIDATES = [
    Path(r"C:\Windows\Fonts\msyh.ttc"),
    Path(r"C:\Windows\Fonts\simsun.ttc"),
    Path(r"C:\Windows\Fonts\arial.ttf"),
]
PAGE_SIZE = (1240, 1754)
PDF_PAGE = (595, 842)
PDF_MARGIN = 54
REPORT_TITLE = "华东智造集团 2026 年半年度授信风险监测报告"
REPORT_META = "报告日期：2026-06-30    资料库：企业授信与合规审查    文档编号：FR-CR-2026-H1-017"
REPORT_SECTIONS = [
    {
        "title": "一、执行摘要",
        "paragraphs": [
            "华东智造集团主营工业自动化设备、精密零部件和售后维保服务。2026 年上半年，公司订单规模保持增长，"
            "但下游新能源客户账期拉长，应收账款周转效率下降。本报告用于授信复核、贷后监测和风险预警，"
            "重点关注现金流覆盖能力、客户集中度、短债偿付弹性和存货跌价风险。",
        ],
    },
    {
        "title": "二、经营与财务指标",
        "paragraphs": [
            "经营指标显示，公司收入端仍具备韧性，但经营现金流同比下降明显。授信审查应结合核心客户回款计划、"
            "在手订单质量和短期债务到期结构进行综合判断。",
        ],
        "table": {
            "headers": ["指标", "2025H1", "2026H1", "同比变化", "风险提示"],
            "rows": [
                ["营业收入", "12.8 亿元", "15.4 亿元", "+20.3%", "订单增长较快"],
                ["毛利率", "31.2%", "29.6%", "-1.6pct", "原材料价格上行"],
                ["经营现金流", "1.9 亿元", "0.8 亿元", "-57.9%", "回款周期延长"],
                ["应收账款周转天数", "74 天", "96 天", "+22 天", "赊销压力上升"],
                ["存货周转天数", "88 天", "104 天", "+16 天", "备货周期拉长"],
            ],
        },
    },
    {
        "title": "三、授信风险矩阵",
        "paragraphs": [
            "风险矩阵中，客户集中度和逾期应收款占比已触发关注阈值。建议将相关指标纳入贷后月度监测，"
            "并要求企业按月提交应收账款账龄、回款责任人和重点客户确认函。",
        ],
        "table": {
            "headers": ["风险项", "当前状态", "阈值", "评级影响", "处置建议"],
            "rows": [
                ["资产负债率", "64.8%", "不高于 70%", "稳定", "维持季度跟踪"],
                ["前五大客户收入占比", "48.5%", "不高于 45%", "关注", "补充客户分散计划"],
                ["逾期应收款占比", "6.2%", "不高于 5%", "关注", "要求回款专项台账"],
                ["短债覆盖倍数", "1.18x", "不低于 1.20x", "轻微承压", "保留提款节奏限制"],
            ],
        },
    },
    {
        "title": "四、现金流压力测试",
        "paragraphs": [
            "在基准情景下，公司仍可覆盖未来十二个月短期债务；在压力情景下，若回款延迟超过 45 天，"
            "短债覆盖倍数将降至 1.05x 以下。银行应保留提款节奏限制，并对新增敞口设置资金用途核验。",
        ],
        "table": {
            "headers": ["情景", "回款假设", "经营现金流", "短债覆盖倍数", "授信建议"],
            "rows": [
                ["基准", "主要客户按合同回款", "1.6 亿元", "1.24x", "维持额度"],
                ["轻度压力", "回款平均延迟 30 天", "0.9 亿元", "1.11x", "限制新增提款"],
                ["重度压力", "回款平均延迟 45 天", "0.3 亿元", "0.96x", "暂停新增敞口"],
            ],
        },
    },
    {
        "title": "五、整改事项跟踪",
        "table": {
            "headers": ["事项", "责任部门", "到期时间", "当前进展", "贷后要求"],
            "rows": [
                ["客户集中度压降", "销售中心", "2026-09-30", "已提交替代客户名单", "每月更新订单转化"],
                ["逾期应收款清收", "财务中心", "2026-08-31", "完成 42% 回款", "提供回款流水"],
                ["存货库龄治理", "供应链中心", "2026-10-15", "长库龄物料重分类", "提交跌价测试"],
            ],
        },
    },
    {
        "title": "六、审查结论",
        "paragraphs": [
            "建议维持综合授信额度 3.5 亿元，其中流动资金贷款 2.2 亿元、银行承兑汇票 0.8 亿元、"
            "保函额度 0.5 亿元。新增提款需满足两个条件：一是 2026 年三季度经营现金流转正；"
            "二是逾期应收款占比降至 5% 以下。",
        ],
        "bullets": [
            "将前五大客户集中度纳入贷后月度监测。",
            "要求企业提供核心客户回款计划和存货库龄明细。",
            "若短债覆盖倍数连续两个月低于 1.1x，应暂停新增敞口。",
        ],
    },
]


def parse_args() -> argparse.Namespace:
    """解析命令行参数"""
    parser = argparse.ArgumentParser(description="生成一组多格式企业金融文档，用于演示 Docling 解析和分块")
    parser.add_argument("--out-dir", type=Path, default=DEFAULT_OUTPUT_DIR, help="输出文档目录")
    parser.add_argument(
        "--pdf-out",
        "--out",
        dest="pdf_out",
        type=Path,
        default=None,
        help="可选 PDF 输出路径；未提供时写入 out-dir",
    )
    return parser.parse_args()


def main() -> int:
    """生成多格式演示文件。"""
    args = parse_args()
    output_dir = args.out_dir.resolve()
    pdf_output = args.pdf_out.resolve() if args.pdf_out else None
    generated = generate_demo_documents(output_dir, pdf_output=pdf_output)
    for path in generated:
        print(path)
    return 0


def generate_demo_documents(output_dir: Path, *, pdf_output: Path | None = None) -> list[Path]:
    """生成企业 RAG 常见格式的演示文件。"""
    output_dir.mkdir(parents=True, exist_ok=True)
    pdf_path = pdf_output or output_dir / f"{DEFAULT_BASENAME}.pdf"
    pdf_path.parent.mkdir(parents=True, exist_ok=True)
    save_demo_pdf(pdf_path)

    markdown_path = write_markdown(output_dir / f"{DEFAULT_BASENAME}.md")
    paths = [
        pdf_path,
        markdown_path,
        write_text(output_dir / f"{DEFAULT_BASENAME}.txt"),
        write_csv(output_dir / f"{DEFAULT_BASENAME}.csv"),
        write_docling_json(markdown_path, output_dir / f"{DEFAULT_BASENAME}.json"),
        write_html(output_dir / f"{DEFAULT_BASENAME}.html"),
        write_docx(output_dir / f"{DEFAULT_BASENAME}.docx"),
        write_xlsx(output_dir / f"{DEFAULT_BASENAME}.xlsx"),
        write_pptx(output_dir / f"{DEFAULT_BASENAME}.pptx"),
        write_scan_image(output_dir / f"{DEFAULT_BASENAME}_scan.png"),
    ]
    return paths


def save_demo_pdf(output: Path) -> None:
    """生成 PDF 文件。"""
    font_path = next((path for path in FONT_CANDIDATES if path.exists()), None)
    if font_path is None:
        raise RuntimeError("未找到可用字体，无法生成中文演示 PDF")

    TextPdfRenderer(font_path).render(output)


def write_markdown(path: Path) -> Path:
    """写入 Markdown 演示文档。"""
    lines = [f"# {REPORT_TITLE}", "", REPORT_META, ""]
    for section in REPORT_SECTIONS:
        lines.extend([f"## {section['title']}", ""])
        for paragraph in section.get("paragraphs", []):
            lines.extend([paragraph, ""])
        if "table" in section:
            table = section["table"]
            lines.extend(markdown_table(table["headers"], table["rows"]))
            lines.append("")
        for bullet in section.get("bullets", []):
            lines.append(f"- {bullet}")
        if section.get("bullets"):
            lines.append("")
    path.write_text("\n".join(lines).rstrip() + "\n", encoding="utf-8")
    return path


def write_text(path: Path) -> Path:
    """写入纯文本演示文档。"""
    lines = [REPORT_TITLE, REPORT_META, ""]
    for section in REPORT_SECTIONS:
        lines.extend([section["title"], ""])
        lines.extend(section.get("paragraphs", []))
        if "table" in section:
            table = section["table"]
            lines.append("\t".join(table["headers"]))
            lines.extend("\t".join(row) for row in table["rows"])
        lines.extend(f"- {bullet}" for bullet in section.get("bullets", []))
        lines.append("")
    path.write_text("\n".join(lines).rstrip() + "\n", encoding="utf-8")
    return path


def write_csv(path: Path) -> Path:
    """写入包含全部表格行的 CSV 演示文档。"""
    with path.open("w", encoding="utf-8-sig", newline="") as file:
        writer = csv.writer(file)
        writer.writerow(["section", "table_column_1", "table_column_2", "table_column_3", "table_column_4", "table_column_5"])
        for section in REPORT_SECTIONS:
            table = section.get("table")
            if not table:
                continue
            writer.writerow([section["title"], *table["headers"]])
            for row in table["rows"]:
                writer.writerow([section["title"], *row])
    return path


def write_docling_json(markdown_path: Path, path: Path) -> Path:
    """将 Markdown 样例转换为 DoclingDocument JSON，匹配 Docling 的 .json 输入语义。"""
    try:
        from llama_index.readers.docling import DoclingReader
    except Exception as exc:
        raise RuntimeError("生成 Docling JSON 样例需要安装 llama-index-readers-docling 依赖") from exc

    export_type = getattr(getattr(DoclingReader, "ExportType", object), "JSON", "json")
    reader = DoclingReader(export_type=export_type)
    documents = reader.load_data(str(markdown_path))
    for document in documents:
        text = document_text(document).strip()
        if not text:
            continue
        path.write_text(json.dumps(json.loads(text), ensure_ascii=False, indent=2), encoding="utf-8")
        return path
    raise RuntimeError(f"Docling 未能从 Markdown 样例生成 JSON: {markdown_path}")


def document_text(document: object) -> str:
    """读取 LlamaIndex Document 文本。"""
    text = getattr(document, "text", None)
    if text is not None:
        return str(text or "")
    getter = getattr(document, "get_content", None)
    if callable(getter):
        return str(getter() or "")
    return str(document or "")


def write_html(path: Path) -> Path:
    """写入 HTML 演示文档。"""
    parts = [
        "<!doctype html>",
        '<html lang="zh-CN">',
        "<head>",
        '<meta charset="utf-8">',
        f"<title>{escape(REPORT_TITLE)}</title>",
        "</head>",
        "<body>",
        f"<h1>{escape(REPORT_TITLE)}</h1>",
        f"<p>{escape(REPORT_META)}</p>",
    ]
    for section in REPORT_SECTIONS:
        parts.append(f"<h2>{escape(section['title'])}</h2>")
        for paragraph in section.get("paragraphs", []):
            parts.append(f"<p>{escape(paragraph)}</p>")
        if "table" in section:
            table = section["table"]
            parts.append("<table>")
            parts.append("<thead><tr>" + "".join(f"<th>{escape(value)}</th>" for value in table["headers"]) + "</tr></thead>")
            parts.append("<tbody>")
            for row in table["rows"]:
                parts.append("<tr>" + "".join(f"<td>{escape(value)}</td>" for value in row) + "</tr>")
            parts.append("</tbody></table>")
        if section.get("bullets"):
            parts.append("<ul>")
            parts.extend(f"<li>{escape(bullet)}</li>" for bullet in section["bullets"])
            parts.append("</ul>")
    parts.extend(["</body>", "</html>"])
    path.write_text("\n".join(parts), encoding="utf-8")
    return path


def write_docx(path: Path) -> Path:
    """写入一个不依赖 python-docx 的简洁 DOCX。"""
    document_xml = "\n".join(
        [
            '<?xml version="1.0" encoding="UTF-8" standalone="yes"?>',
            '<w:document xmlns:w="http://schemas.openxmlformats.org/wordprocessingml/2006/main">',
            "<w:body>",
            docx_paragraph(REPORT_TITLE, style="Title"),
            docx_paragraph(REPORT_META),
            *docx_sections(),
            '<w:sectPr><w:pgSz w:w="11906" w:h="16838"/><w:pgMar w:top="1440" w:right="1440" w:bottom="1440" w:left="1440"/></w:sectPr>',
            "</w:body>",
            "</w:document>",
        ]
    )
    with zipfile.ZipFile(path, "w", compression=zipfile.ZIP_DEFLATED) as docx:
        docx.writestr(
            "[Content_Types].xml",
            """<?xml version="1.0" encoding="UTF-8" standalone="yes"?>
<Types xmlns="http://schemas.openxmlformats.org/package/2006/content-types">
  <Default Extension="rels" ContentType="application/vnd.openxmlformats-package.relationships+xml"/>
  <Default Extension="xml" ContentType="application/xml"/>
  <Override PartName="/word/document.xml" ContentType="application/vnd.openxmlformats-officedocument.wordprocessingml.document.main+xml"/>
  <Override PartName="/word/styles.xml" ContentType="application/vnd.openxmlformats-officedocument.wordprocessingml.styles+xml"/>
</Types>""",
        )
        docx.writestr(
            "_rels/.rels",
            """<?xml version="1.0" encoding="UTF-8" standalone="yes"?>
<Relationships xmlns="http://schemas.openxmlformats.org/package/2006/relationships">
  <Relationship Id="rId1" Type="http://schemas.openxmlformats.org/officeDocument/2006/relationships/officeDocument" Target="word/document.xml"/>
</Relationships>""",
        )
        docx.writestr(
            "word/_rels/document.xml.rels",
            """<?xml version="1.0" encoding="UTF-8" standalone="yes"?>
<Relationships xmlns="http://schemas.openxmlformats.org/package/2006/relationships">
  <Relationship Id="rIdStyles" Type="http://schemas.openxmlformats.org/officeDocument/2006/relationships/styles" Target="styles.xml"/>
</Relationships>""",
        )
        docx.writestr("word/styles.xml", docx_styles_xml())
        docx.writestr("word/document.xml", document_xml)
    return path


def write_xlsx(path: Path) -> Path:
    """写入带格式的 XLSX 指标台账。"""
    try:
        from openpyxl import Workbook
        from openpyxl.styles import Alignment, Border, Font, PatternFill, Side
        from openpyxl.utils import get_column_letter
    except Exception as exc:
        raise RuntimeError("生成 XLSX 样例需要安装 openpyxl") from exc

    workbook = Workbook()
    workbook.remove(workbook.active)
    header_fill = PatternFill("solid", fgColor="E8F0F6")
    thin_border = Border(
        left=Side(style="thin", color="B8C2CC"),
        right=Side(style="thin", color="B8C2CC"),
        top=Side(style="thin", color="B8C2CC"),
        bottom=Side(style="thin", color="B8C2CC"),
    )

    def style_table(sheet, header_row: int, max_row: int, max_col: int) -> None:
        for cell in sheet[header_row]:
            cell.fill = header_fill
            cell.font = Font(name="Microsoft YaHei", bold=True, color="111827")
            cell.alignment = Alignment(horizontal="center", vertical="center", wrap_text=True)
        for row in sheet.iter_rows(min_row=header_row, max_row=max_row, max_col=max_col):
            for cell in row:
                cell.border = thin_border
                cell.alignment = Alignment(vertical="center", wrap_text=True)
                font = copy(cell.font)
                font.name = "Microsoft YaHei"
                cell.font = font
        sheet.freeze_panes = f"A{header_row + 1}"

    metrics = workbook.create_sheet("经营指标")
    workbook.properties.title = REPORT_TITLE
    metrics.append(["指标", "2025H1", "2026H1", "同比变化", "风险提示"])
    metric_rows = [
        ["营业收入(亿元)", 12.8, 15.4, 0.203, "订单增长较快"],
        ["毛利率", 0.312, 0.296, -0.016, "原材料价格上行"],
        ["经营现金流(亿元)", 1.9, 0.8, -0.579, "回款周期延长"],
        ["应收账款周转天数", 74, 96, 22, "赊销压力上升"],
        ["存货周转天数", 88, 104, 16, "备货周期拉长"],
    ]
    for row in metric_rows:
        metrics.append(row)
    for row in metrics.iter_rows(min_row=2, max_row=6, min_col=2, max_col=4):
        for cell in row:
            if cell.row == 3 or cell.column == 4 and cell.row in {2, 3, 4}:
                cell.number_format = "0.0%"
            elif cell.row in {5, 6}:
                cell.number_format = "#,##0"
            else:
                cell.number_format = "#,##0.0"
    style_table(metrics, 1, 6, 5)

    risk = workbook.create_sheet("授信风险矩阵")
    risk.append(["风险项", "当前状态", "阈值", "评级影响", "处置建议"])
    for row in REPORT_SECTIONS[2]["table"]["rows"]:
        risk.append(row)
    style_table(risk, 1, 5, 5)

    actions = workbook.create_sheet("整改跟踪")
    actions.append(["事项", "责任部门", "到期时间", "当前进展", "贷后要求", "状态"])
    for row in REPORT_SECTIONS[4]["table"]["rows"]:
        actions.append([*row, "进行中"])
    style_table(actions, 1, 4, 6)

    for sheet in workbook.worksheets:
        for width_index, width in enumerate([20, 18, 18, 18, 26, 12], start=1):
            sheet.column_dimensions[get_column_letter(width_index)].width = width
        for row in sheet.iter_rows():
            sheet.row_dimensions[row[0].row].height = 24

    workbook.save(path)
    return path


def write_pptx(path: Path) -> Path:
    """写入带排版的 PPTX 授信审查简报。"""
    try:
        from pptx import Presentation
        from pptx.dml.color import RGBColor
        from pptx.enum.text import PP_ALIGN
        from pptx.util import Inches, Pt
    except Exception as exc:
        raise RuntimeError("生成 PPTX 样例需要安装 python-pptx") from exc

    presentation = Presentation()
    presentation.slide_width = Inches(13.333)
    presentation.slide_height = Inches(7.5)

    def add_title(slide, title: str, subtitle: str | None = None) -> None:
        shape = slide.shapes.add_textbox(Inches(0.55), Inches(0.35), Inches(12.2), Inches(0.55))
        paragraph = shape.text_frame.paragraphs[0]
        paragraph.text = title
        paragraph.font.name = "Microsoft YaHei"
        paragraph.font.bold = True
        paragraph.font.size = Pt(24)
        paragraph.font.color.rgb = RGBColor(18, 50, 74)
        if subtitle:
            sub = slide.shapes.add_textbox(Inches(0.58), Inches(0.95), Inches(12.0), Inches(0.35))
            paragraph = sub.text_frame.paragraphs[0]
            paragraph.text = subtitle
            paragraph.font.name = "Microsoft YaHei"
            paragraph.font.size = Pt(10)
            paragraph.font.color.rgb = RGBColor(92, 101, 112)

    def add_bullets(slide, items: list[str], x: float, y: float, width: float, height: float) -> None:
        shape = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(width), Inches(height))
        frame = shape.text_frame
        frame.clear()
        for index, item in enumerate(items):
            paragraph = frame.paragraphs[0] if index == 0 else frame.add_paragraph()
            paragraph.text = item
            paragraph.level = 0
            paragraph.font.name = "Microsoft YaHei"
            paragraph.font.size = Pt(14)
            paragraph.font.color.rgb = RGBColor(31, 41, 51)

    def add_table(slide, headers: list[str], rows: list[list[str]], x: float, y: float, width: float, height: float) -> None:
        table = slide.shapes.add_table(len(rows) + 1, len(headers), Inches(x), Inches(y), Inches(width), Inches(height)).table
        for column_index in range(len(headers)):
            table.columns[column_index].width = int(Inches(width / len(headers)))
        for column_index, header in enumerate(headers):
            cell = table.cell(0, column_index)
            cell.text = header
            cell.fill.solid()
            cell.fill.fore_color.rgb = RGBColor(232, 240, 246)
            cell.text_frame.paragraphs[0].font.bold = True
        for row_index, row in enumerate(rows, start=1):
            for column_index, value in enumerate(row):
                cell = table.cell(row_index, column_index)
                cell.text = value
                for paragraph in cell.text_frame.paragraphs:
                    paragraph.font.name = "Microsoft YaHei"
                    paragraph.font.size = Pt(8)
                    paragraph.alignment = PP_ALIGN.CENTER if column_index in {1, 2, 3} else PP_ALIGN.LEFT

    slide = presentation.slides.add_slide(presentation.slide_layouts[6])
    add_title(slide, REPORT_TITLE, REPORT_META)
    add_bullets(
        slide,
        [
            "2026H1 营业收入 15.4 亿元，同比增长 20.3%，但毛利率下降至 29.6%。",
            "经营现金流降至 0.8 亿元，应收账款周转天数延长至 96 天。",
            "建议维持授信额度 3.5 亿元，同时设置提款节奏和贷后月度监测条件。",
        ],
        0.75,
        1.55,
        6.0,
        2.2,
    )
    add_table(slide, REPORT_SECTIONS[1]["table"]["headers"], REPORT_SECTIONS[1]["table"]["rows"][:4], 0.75, 4.0, 11.8, 2.2)

    slide = presentation.slides.add_slide(presentation.slide_layouts[6])
    add_title(slide, "授信风险监测矩阵", "客户集中度、逾期应收款和短债覆盖倍数是本期核心关注项")
    add_table(slide, REPORT_SECTIONS[2]["table"]["headers"], REPORT_SECTIONS[2]["table"]["rows"], 0.75, 1.35, 11.8, 3.0)
    add_bullets(
        slide,
        ["处置重点：补充客户分散计划、建立回款专项台账、保留提款节奏限制。"],
        0.85,
        5.0,
        11.5,
        0.8,
    )

    slide = presentation.slides.add_slide(presentation.slide_layouts[6])
    add_title(slide, "审查结论与后续动作", "以现金流改善和逾期应收款压降作为新增提款前置条件")
    add_table(slide, REPORT_SECTIONS[3]["table"]["headers"], REPORT_SECTIONS[3]["table"]["rows"], 0.75, 1.35, 11.8, 2.5)
    add_bullets(slide, REPORT_SECTIONS[5]["bullets"], 0.85, 4.35, 11.5, 1.6)

    presentation.save(path)
    return path


def write_scan_image(path: Path) -> Path:
    """写入图片式 OCR 样例。"""
    font_path = next((candidate for candidate in FONT_CANDIDATES if candidate.exists()), None)
    if font_path is None:
        raise RuntimeError("未找到可用字体，无法生成 OCR 图片样例")
    ReportRenderer(font_path).page_one().save(path, "PNG")
    return path


def docx_styles_xml() -> str:
    """生成最小但真实的 Word 样式表，让 Docling 能识别标题层级。"""
    return """<?xml version="1.0" encoding="UTF-8" standalone="yes"?>
<w:styles xmlns:w="http://schemas.openxmlformats.org/wordprocessingml/2006/main">
  <w:style w:type="paragraph" w:default="1" w:styleId="Normal">
    <w:name w:val="Normal"/>
    <w:qFormat/>
    <w:pPr>
      <w:spacing w:after="120" w:line="276" w:lineRule="auto"/>
    </w:pPr>
    <w:rPr>
      <w:rFonts w:ascii="Microsoft YaHei" w:eastAsia="Microsoft YaHei" w:hAnsi="Microsoft YaHei"/>
      <w:sz w:val="21"/>
    </w:rPr>
  </w:style>
  <w:style w:type="paragraph" w:styleId="Title">
    <w:name w:val="Title"/>
    <w:basedOn w:val="Normal"/>
    <w:next w:val="Normal"/>
    <w:qFormat/>
    <w:pPr>
      <w:spacing w:before="120" w:after="180"/>
      <w:jc w:val="center"/>
    </w:pPr>
    <w:rPr>
      <w:b/>
      <w:rFonts w:ascii="Microsoft YaHei" w:eastAsia="Microsoft YaHei" w:hAnsi="Microsoft YaHei"/>
      <w:color w:val="12324A"/>
      <w:sz w:val="32"/>
    </w:rPr>
  </w:style>
  <w:style w:type="paragraph" w:styleId="Heading1">
    <w:name w:val="heading 1"/>
    <w:basedOn w:val="Normal"/>
    <w:next w:val="Normal"/>
    <w:uiPriority w:val="9"/>
    <w:qFormat/>
    <w:pPr>
      <w:keepNext/>
      <w:spacing w:before="240" w:after="120"/>
      <w:outlineLvl w:val="0"/>
    </w:pPr>
    <w:rPr>
      <w:b/>
      <w:rFonts w:ascii="Microsoft YaHei" w:eastAsia="Microsoft YaHei" w:hAnsi="Microsoft YaHei"/>
      <w:color w:val="0F3B36"/>
      <w:sz w:val="27"/>
    </w:rPr>
  </w:style>
</w:styles>"""


def docx_sections() -> list[str]:
    """转换报告段落和表格为 DOCX XML 片段。"""
    parts: list[str] = []
    for section in REPORT_SECTIONS:
        parts.append(docx_paragraph(section["title"], style="Heading1"))
        parts.extend(docx_paragraph(paragraph) for paragraph in section.get("paragraphs", []))
        if "table" in section:
            table = section["table"]
            parts.append(docx_table(table["headers"], table["rows"]))
        parts.extend(docx_paragraph(f"• {bullet}") for bullet in section.get("bullets", []))
    return parts


def docx_paragraph(text: str, *, style: str | None = None, bold: bool = False) -> str:
    """生成 DOCX 段落 XML。"""
    style_xml = f'<w:pPr><w:pStyle w:val="{style}"/></w:pPr>' if style else ""
    bold_xml = "<w:rPr><w:b/></w:rPr>" if bold else ""
    return f"<w:p>{style_xml}<w:r>{bold_xml}<w:t>{escape(text)}</w:t></w:r></w:p>"


def docx_table(headers: list[str], rows: list[list[str]]) -> str:
    """生成 DOCX 表格 XML。"""
    widths = [1800, 1500, 1500, 1500, 2760]
    table_rows = [
        docx_table_row(headers, widths, is_header=True),
        *[docx_table_row(row, widths, is_header=False) for row in rows],
    ]
    grid_columns = "".join(f'<w:gridCol w:w="{width}"/>' for width in widths)
    return (
        "<w:tbl>"
        "<w:tblPr>"
        '<w:tblW w:w="9060" w:type="dxa"/>'
        '<w:tblBorders><w:top w:val="single" w:sz="4" w:color="B8C2CC"/>'
        '<w:left w:val="single" w:sz="4" w:color="B8C2CC"/>'
        '<w:bottom w:val="single" w:sz="4" w:color="B8C2CC"/>'
        '<w:right w:val="single" w:sz="4" w:color="B8C2CC"/>'
        '<w:insideH w:val="single" w:sz="4" w:color="B8C2CC"/>'
        '<w:insideV w:val="single" w:sz="4" w:color="B8C2CC"/></w:tblBorders>'
        '<w:tblCellMar><w:top w:w="80" w:type="dxa"/><w:left w:w="100" w:type="dxa"/>'
        '<w:bottom w:w="80" w:type="dxa"/><w:right w:w="100" w:type="dxa"/></w:tblCellMar>'
        "</w:tblPr>"
        f"<w:tblGrid>{grid_columns}</w:tblGrid>"
        + "".join(table_rows)
        + "</w:tbl>"
    )


def docx_table_row(values: list[str], widths: list[int], *, is_header: bool) -> str:
    """生成 DOCX 表格行 XML。"""
    header_props = "<w:trPr><w:tblHeader/></w:trPr>" if is_header else ""
    shading = '<w:shd w:fill="E8F0F6"/>' if is_header else ""
    cells = "".join(
        f'<w:tc><w:tcPr><w:tcW w:w="{width}" w:type="dxa"/>'
        f"{shading}"
        f"</w:tcPr>{docx_paragraph(value, bold=is_header)}</w:tc>"
        for value, width in zip(values, widths)
    )
    return f"<w:tr>{header_props}{cells}</w:tr>"


def markdown_table(headers: list[str], rows: list[list[str]]) -> list[str]:
    """生成 Markdown 表格。"""
    return [
        "| " + " | ".join(headers) + " |",
        "| " + " | ".join("---" for _ in headers) + " |",
        *["| " + " | ".join(row) + " |" for row in rows],
    ]


class TextPdfRenderer:
    """使用真实文本绘制 PDF，避免 OCR 型样例的不稳定性。"""

    def __init__(self, font_path: Path):
        self.font_path = font_path
        self.doc = fitz.open()
        self.page: fitz.Page
        self.page_no = 0
        self.y = 0.0
        self.font_title = ImageFont.truetype(str(font_path), 18)
        self.font_h1 = ImageFont.truetype(str(font_path), 14)
        self.font_body = ImageFont.truetype(str(font_path), 10)
        self.font_small = ImageFont.truetype(str(font_path), 8)
        self.font_table = ImageFont.truetype(str(font_path), 8)

    def render(self, output: Path) -> None:
        """渲染并保存 PDF。"""
        self.new_page()
        self.title(REPORT_TITLE)
        self.meta(REPORT_META)
        for section in REPORT_SECTIONS:
            self.heading(section["title"])
            for paragraph in section.get("paragraphs", []):
                self.paragraph(paragraph)
            if "table" in section:
                table = section["table"]
                self.table(table["headers"], table["rows"])
            for bullet in section.get("bullets", []):
                self.bullet(bullet)
        self.doc.save(output, garbage=4, deflate=True)
        self.doc.close()

    def new_page(self) -> None:
        """创建带页眉页脚的新页面。"""
        self.page_no += 1
        self.page = self.doc.new_page(width=PDF_PAGE[0], height=PDF_PAGE[1])
        self.y = 86
        self.text(PDF_MARGIN, 42, "企业授信与合规审查资料库", size=8, color=(0.34, 0.38, 0.42))
        self.text(410, 42, "内部资料 - 风险监测样例", size=8, color=(0.34, 0.38, 0.42))
        self.page.draw_line((PDF_MARGIN, 56), (PDF_PAGE[0] - PDF_MARGIN, 56), color=(0.82, 0.85, 0.88), width=0.8)
        self.text(286, 814, f"第 {self.page_no} 页", size=8, color=(0.42, 0.45, 0.50))

    def title(self, value: str) -> None:
        """写入标题。"""
        self.ensure_space(52)
        self.text(PDF_MARGIN, self.y, value, size=18, color=(0.07, 0.20, 0.29))
        self.y += 30
        self.page.draw_line((PDF_MARGIN, self.y), (PDF_PAGE[0] - PDF_MARGIN, self.y), color=(0.66, 0.72, 0.78), width=0.8)
        self.y += 18

    def meta(self, value: str) -> None:
        """写入报告元信息。"""
        self.ensure_space(24)
        self.text(PDF_MARGIN, self.y, value, size=8, color=(0.41, 0.44, 0.49))
        self.y += 26

    def heading(self, value: str) -> None:
        """写入章节标题。"""
        self.ensure_space(44)
        self.y += 8
        self.text(PDF_MARGIN, self.y, value, size=14, color=(0.06, 0.23, 0.21))
        self.y += 26

    def paragraph(self, value: str) -> None:
        """写入段落。"""
        max_width = PDF_PAGE[0] - PDF_MARGIN * 2
        for line in wrap_text(value, self.font_body, int(max_width)):
            self.ensure_space(16)
            self.text(PDF_MARGIN, self.y, line, size=10, color=(0.12, 0.16, 0.20))
            self.y += 15
        self.y += 8

    def bullet(self, value: str) -> None:
        """写入项目符号。"""
        max_width = PDF_PAGE[0] - PDF_MARGIN * 2 - 16
        for index, line in enumerate(wrap_text(value, self.font_body, int(max_width))):
            self.ensure_space(16)
            prefix = "- " if index == 0 else "  "
            self.text(PDF_MARGIN + 12, self.y, f"{prefix}{line}", size=10, color=(0.12, 0.16, 0.20))
            self.y += 15

    def table(self, headers: list[str], rows: list[list[str]]) -> None:
        """写入简单表格，保留文本和网格线。"""
        widths = [92, 84, 84, 88, 135]
        row_padding = 5
        all_rows = [headers, *rows]
        wrapped_rows: list[list[list[str]]] = []
        row_heights: list[int] = []
        for row in all_rows:
            wrapped_cells = [
                wrap_text(cell, self.font_table, max(int(width * 1.25), 20))
                for cell, width in zip(row, widths)
            ]
            wrapped_rows.append(wrapped_cells)
            row_heights.append(max(24, max(len(lines) for lines in wrapped_cells) * 10 + row_padding * 2))

        total_height = sum(row_heights) + 12
        self.ensure_space(total_height)
        x0 = PDF_MARGIN
        y = self.y
        for row_index, (row, height) in enumerate(zip(wrapped_rows, row_heights)):
            fill = (0.91, 0.94, 0.96) if row_index == 0 else ((1, 1, 1) if row_index % 2 else (0.97, 0.98, 0.99))
            self.page.draw_rect(
                fitz.Rect(x0, y, x0 + sum(widths), y + height),
                color=(0.72, 0.76, 0.80),
                fill=fill,
                width=0.5,
            )
            x = x0
            for cell_lines, width in zip(row, widths):
                self.page.draw_line((x, y), (x, y + height), color=(0.72, 0.76, 0.80), width=0.5)
                text_y = y + row_padding + 8
                for line in cell_lines:
                    self.text(x + 4, text_y, line, size=8, color=(0.07, 0.09, 0.12))
                    text_y += 10
                x += width
            self.page.draw_line((x, y), (x, y + height), color=(0.72, 0.76, 0.80), width=0.5)
            y += height
        self.y = y + 14

    def text(self, x: float, y: float, value: str, *, size: int, color: tuple[float, float, float]) -> None:
        """写入一行可抽取文本。"""
        self.page.insert_text(
            (x, y),
            value,
            fontsize=size,
            fontname="finragdemo",
            fontfile=str(self.font_path),
            color=color,
        )

    def ensure_space(self, needed: float) -> None:
        """空间不足时换页。"""
        if self.y + needed > PDF_PAGE[1] - 56:
            self.new_page()


class ReportRenderer:
    """企业报告图片渲染器"""

    def __init__(self, font_path: Path):
        self.font_path = font_path
        self.font_title = ImageFont.truetype(str(font_path), 34)
        self.font_h1 = ImageFont.truetype(str(font_path), 25)
        self.font_body = ImageFont.truetype(str(font_path), 20)
        self.font_small = ImageFont.truetype(str(font_path), 16)
        self.font_table = ImageFont.truetype(str(font_path), 17)
        self.page: Image.Image
        self.draw: ImageDraw.ImageDraw
        self.y = 0

    def render_pages(self) -> list[Image.Image]:
        """渲染所有页面"""
        return [self.page_one(), self.page_two()]

    def page_one(self) -> Image.Image:
        """渲染第一页"""
        self.new_page(1)
        self.title("华东智造集团 2026 年半年度授信风险监测报告")
        self.meta("报告日期：2026-06-30    资料库：企业授信与合规审查    文档编号：FR-CR-2026-H1-017")
        self.heading("一、执行摘要")
        self.paragraph(
            "华东智造集团主营工业自动化设备、精密零部件和售后维保服务。2026 年上半年，公司订单规模保持增长，"
            "但下游新能源客户账期拉长，应收账款周转效率下降。本报告用于授信复核、贷后监测和风险预警，"
            "重点关注现金流覆盖能力、客户集中度、短债偿付弹性和存货跌价风险。"
        )
        self.heading("二、经营与财务指标")
        self.table(
            ["指标", "2025H1", "2026H1", "同比变化", "风险提示"],
            [
                ["营业收入", "12.8 亿元", "15.4 亿元", "+20.3%", "订单增长较快"],
                ["毛利率", "31.2%", "29.6%", "-1.6pct", "原材料价格上行"],
                ["经营现金流", "1.9 亿元", "0.8 亿元", "-57.9%", "回款周期延长"],
                ["应收账款周转天数", "74 天", "96 天", "+22 天", "赊销压力上升"],
                ["存货周转天数", "88 天", "104 天", "+16 天", "备货周期拉长"],
            ],
            [170, 150, 150, 150, 300],
        )
        self.paragraph(
            "经营指标显示，公司收入端仍具备韧性，但经营现金流同比下降明显。授信审查应结合核心客户回款计划、"
            "在手订单质量和短期债务到期结构进行综合判断。"
        )
        self.heading("三、授信风险矩阵")
        self.table(
            ["风险项", "当前状态", "阈值", "评级影响", "处置建议"],
            [
                ["资产负债率", "64.8%", "不高于 70%", "稳定", "维持季度跟踪"],
                ["前五大客户收入占比", "48.5%", "不高于 45%", "关注", "补充客户分散计划"],
                ["逾期应收款占比", "6.2%", "不高于 5%", "关注", "要求回款专项台账"],
                ["短债覆盖倍数", "1.18x", "不低于 1.20x", "轻微承压", "保留提款节奏限制"],
            ],
            [190, 165, 165, 165, 300],
        )
        self.paragraph(
            "风险矩阵中，客户集中度和逾期应收款占比已触发关注阈值。建议将相关指标纳入贷后月度监测，"
            "并要求企业按月提交应收账款账龄、回款责任人和重点客户确认函。"
        )
        return self.page

    def page_two(self) -> Image.Image:
        """渲染第二页"""
        self.new_page(2)
        self.heading("四、现金流压力测试")
        self.table(
            ["情景", "回款假设", "经营现金流", "短债覆盖倍数", "授信建议"],
            [
                ["基准", "主要客户按合同回款", "1.6 亿元", "1.24x", "维持额度"],
                ["轻度压力", "回款平均延迟 30 天", "0.9 亿元", "1.11x", "限制新增提款"],
                ["重度压力", "回款平均延迟 45 天", "0.3 亿元", "0.96x", "暂停新增敞口"],
            ],
            [135, 285, 165, 190, 190],
        )
        self.paragraph(
            "在基准情景下，公司仍可覆盖未来十二个月短期债务；在压力情景下，若回款延迟超过 45 天，"
            "短债覆盖倍数将降至 1.05x 以下。银行应保留提款节奏限制，并对新增敞口设置资金用途核验。"
        )
        self.heading("五、整改事项跟踪")
        self.table(
            ["事项", "责任部门", "到期时间", "当前进展", "贷后要求"],
            [
                ["客户集中度压降", "销售中心", "2026-09-30", "已提交替代客户名单", "每月更新订单转化"],
                ["逾期应收款清收", "财务中心", "2026-08-31", "完成 42% 回款", "提供回款流水"],
                ["存货库龄治理", "供应链中心", "2026-10-15", "长库龄物料重分类", "提交跌价测试"],
            ],
            [190, 150, 160, 250, 230],
        )
        self.heading("六、审查结论")
        self.paragraph(
            "建议维持综合授信额度 3.5 亿元，其中流动资金贷款 2.2 亿元、银行承兑汇票 0.8 亿元、"
            "保函额度 0.5 亿元。新增提款需满足两个条件：一是 2026 年三季度经营现金流转正；"
            "二是逾期应收款占比降至 5% 以下。"
        )
        self.bullets(
            [
                "将前五大客户集中度纳入贷后月度监测。",
                "要求企业提供核心客户回款计划和存货库龄明细。",
                "若短债覆盖倍数连续两个月低于 1.1x，应暂停新增敞口。",
            ]
        )
        return self.page

    def new_page(self, page_no: int) -> None:
        """新建页面"""
        self.page = Image.new("RGB", PAGE_SIZE, "white")
        self.draw = ImageDraw.Draw(self.page)
        self.y = 112
        self.draw.text((92, 48), "企业授信与合规审查资料库", font=self.font_small, fill="#58606A")
        self.draw.text((890, 48), "内部资料 - 风险监测样例", font=self.font_small, fill="#58606A")
        self.draw.line((92, 74, 1148, 74), fill="#D6DCE2", width=2)
        self.draw.text((590, 1708), f"第 {page_no} 页", font=self.font_small, fill="#6B7280", anchor="mm")

    def title(self, text: str) -> None:
        """写入标题"""
        self.draw.text((620, self.y), text, font=self.font_title, fill="#12324A", anchor="ma")
        self.y += 56
        self.draw.line((92, self.y, 1148, self.y), fill="#A9B8C6", width=2)
        self.y += 24

    def meta(self, text: str) -> None:
        """写入报告元信息"""
        self.draw.text((620, self.y), text, font=self.font_small, fill="#68717D", anchor="ma")
        self.y += 42

    def heading(self, text: str) -> None:
        """写入章节标题"""
        self.y += 16
        self.draw.text((92, self.y), text, font=self.font_h1, fill="#0F3B36")
        self.y += 42

    def paragraph(self, text: str) -> None:
        """写入段落"""
        for line in wrap_text(text, self.font_body, 1040):
            self.draw.text((92, self.y), line, font=self.font_body, fill="#1F2933")
            self.y += 34
        self.y += 12

    def bullets(self, items: list[str]) -> None:
        """写入项目符号列表"""
        for item in items:
            for index, line in enumerate(wrap_text(item, self.font_body, 980)):
                prefix = "• " if index == 0 else "  "
                self.draw.text((112, self.y), f"{prefix}{line}", font=self.font_body, fill="#1F2933")
                self.y += 34

    def table(self, headers: list[str], rows: list[list[str]], widths: list[int]) -> None:
        """写入表格"""
        x0 = 92
        row_height = 58
        data = [headers, *rows]
        table_width = sum(widths)
        y = self.y
        for row_index, row in enumerate(data):
            fill = "#E8F0F6" if row_index == 0 else ("#FFFFFF" if row_index % 2 else "#F8FAFC")
            self.draw.rectangle((x0, y, x0 + table_width, y + row_height), outline="#B8C2CC", fill=fill, width=1)
            x = x0
            for value, width in zip(row, widths):
                self.draw.line((x, y, x, y + row_height), fill="#B8C2CC", width=1)
                self.draw_multiline_cell(value, x + 10, y + 12, width - 20)
                x += width
            self.draw.line((x, y, x, y + row_height), fill="#B8C2CC", width=1)
            y += row_height
        self.y = y + 22

    def draw_multiline_cell(self, text: str, x: int, y: int, width: int) -> None:
        """写入单元格文本"""
        lines = wrap_text(text, self.font_table, width)
        for line in lines[:2]:
            self.draw.text((x, y), line, font=self.font_table, fill="#111827")
            y += 23


def wrap_text(text: str, font: ImageFont.FreeTypeFont, max_width: int) -> list[str]:
    """按像素宽度换行"""
    lines: list[str] = []
    current = ""
    for char in text:
        candidate = current + char
        width = font.getlength(candidate)
        if width <= max_width or not current:
            current = candidate
            continue
        lines.append(current)
        current = char
    if current:
        lines.append(current)
    return lines


def save_images_as_pdf(pages: list[Image.Image], output: Path) -> None:
    """将页面图片封装为 PDF"""
    doc = fitz.open()
    with tempfile.TemporaryDirectory() as tmpdir:
        image_paths: list[Path] = []
        for index, page in enumerate(pages, 1):
            image_path = Path(tmpdir) / f"page_{index}.jpg"
            page.save(image_path, "JPEG", quality=88, optimize=True)
            image_paths.append(image_path)
        for image_path in image_paths:
            page = doc.new_page(width=PDF_PAGE[0], height=PDF_PAGE[1])
            page.insert_image(page.rect, filename=str(image_path))
    doc.save(output, garbage=4, deflate=True)
    doc.close()


if __name__ == "__main__":
    raise SystemExit(main())
