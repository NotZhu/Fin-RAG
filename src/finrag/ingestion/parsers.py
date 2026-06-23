"""FinRAG 通用金融文档解析入口"""

from __future__ import annotations

import hashlib
import csv
import io
import json
import logging
import re
from dataclasses import asdict, dataclass
from datetime import datetime, timezone
from html.parser import HTMLParser
from pathlib import Path
from typing import Any, Callable, Dict, Iterable, List, Optional, Sequence

from llama_index.core import Document
from pypdf import PdfReader

logger = logging.getLogger(__name__)

# 支持解析和上传的文件后缀
SUPPORTED_SUFFIXES = {
    ".md",
    ".txt",
    ".pdf",
    ".docx",
    ".csv",
    ".json",
    ".html",
    ".htm",
    ".xlsx",
    ".pptx",
}
# 支持解析的文本文件编码顺序
TEXT_ENCODINGS = ("utf-8-sig", "utf-8", "gb18030", "utf-16", "utf-16-le", "utf-16-be") # 文本文件候选编码顺序


def utc_now_iso() -> str:
    """
    获取当前 UTC 时间的 ISO 字符串
    Returns:
        带时区信息的 ISO 时间字符串
    """
    return datetime.now(timezone.utc).isoformat()


def compute_content_hash(path: Path) -> str:
    """
    计算文件内容 SHA256 哈希，用于文档去重
    Args:
        path: 待计算的文件路径
    Returns:
        sha256: 前缀的内容哈希
    """
    return "sha256:" + hashlib.sha256(path.read_bytes()).hexdigest()


def normalize_text(text: str) -> str:
    """
    规范化文档文本，统一换行、空白和空行
    Args:
        text: 原始文本
    Returns:
        清洗后的文本
    """
    normalized = (text or "").replace("\r\n", "\n").replace("\r", "\n").replace("\u00a0", " ")
    lines: List[str] = []
    blank = False # 上一行是否为空行
    for raw_line in normalized.split("\n"):
        line = re.sub(r"[ \t]+", " ", raw_line).strip() # 统一空白为单个空格并去除行首尾空白
        if line:
            lines.append(line)
            blank = False
        elif not blank: # 当前行为空行且上一行不是空行
            lines.append("") # 保留一个空行
            blank = True
    return "\n".join(lines).strip()


def read_text_file(path: Path) -> str:
    """
    按常见编码顺序读取文本文件
    Args:
        path: 文本文件路径
    Returns:
        读取到的文本；编码无法识别时返回空字符串
    """
    for encoding in TEXT_ENCODINGS:
        try:
            return path.read_text(encoding=encoding)
        except UnicodeError:
            continue
    logger.warning("文本文件编码无法识别: %s", path)
    return ""


def _relative_path(path: Path, data_root: Optional[Path]) -> str:
    """
    计算文件相对数据根目录的路径，失败时返回绝对路径
    Args:
        path: 文件路径
        data_root: 数据根目录
    Returns:
        POSIX 风格路径字符串
    """
    if data_root is not None:
        try:
            return path.resolve().relative_to(data_root.resolve()).as_posix()
        except ValueError:
            pass
    return path.resolve().as_posix()


def is_path_within(path: Path, root: Path) -> bool:
    """
    判断路径是否位于指定根目录内
    Args:
        path: 待判断路径
        root: 根目录
    Returns:
        路径在根目录内时返回 True
    """
    try:
        path.resolve().relative_to(root.resolve())
    except ValueError:
        return False
    return True


def build_base_metadata(
    path: Path,
    text: str,
    file_type: str,
    *,
    knowledge_base_id: str = "default",
    page_number: Optional[int] = None,
    data_root: Optional[Path] = None,
) -> dict:
    """
    为解析出的 Document 构造统一基础元数据
    Args:
        path: 源文件路径
        text: 文档文本
        file_type: 文件类型
        knowledge_base_id: 资料库 ID
        page_number: PDF 页码
        data_root: 数据根目录
    Returns:
        包含 document_id、文件名和页码的元数据字典
    """
    content_hash = compute_content_hash(path)
    relative_path = _relative_path(path, data_root)
    document_seed = f"{knowledge_base_id}:{relative_path}:{content_hash}"
    if page_number is not None:
        document_seed += f":page:{page_number}"
    # 用 MD5 哈希种子计算文档 ID，确保唯一性
    document_id = hashlib.md5(document_seed.encode("utf-8")).hexdigest()
    return {
        "knowledge_base_id": knowledge_base_id,
        "document_id": document_id,
        "source_path": path.resolve().as_posix(),
        "filename": path.name,
        "file_type": file_type,
        "page_number": page_number,
    }


def _markdown_escape(value: Any) -> str:
    """
    转义 Markdown 表格单元格内容
    Args:
        value: 待转义的文本
    Returns:
        转义后的文本，包含 Markdown 特殊字符
    """
    text = normalize_text("" if value is None else str(value))
    return text.replace("|", "\\|").replace("\n", "<br>")


def _markdown_table(rows: Sequence[Sequence[Any]]) -> str:
    """
    将二维表格渲染为 Markdown 表格
    Args:
        rows: 表格行，第一行作为表头
    Returns:
        Markdown 表格文本；空表返回空字符串
    """
    # 转义所有单元格内容
    cleaned_rows = [[_markdown_escape(cell) for cell in row] for row in rows]
    # 过滤掉所有单元格为空的行
    cleaned_rows = [row for row in cleaned_rows if any(cell for cell in row)]
    if not cleaned_rows:
        return ""
    # 计算最大宽度
    width = max(len(row) for row in cleaned_rows)
    # 补充空单元格，使所有行宽度相等
    normalized_rows = [row + [""] * (width - len(row)) for row in cleaned_rows]
    # 提取表头
    header = normalized_rows[0]
    # 提取表格主体
    body = normalized_rows[1:]
    # 构建 Markdown 表格行
    lines = [
        "| " + " | ".join(header) + " |",
        "| " + " | ".join("---" for _ in range(width)) + " |",
    ]
    # 添加表格主体行
    for row in body:
        lines.append("| " + " | ".join(row) + " |")
    return "\n".join(lines)


def _document_from_text(
    path: Path,
    text: str,
    file_type: str,
    *,
    knowledge_base_id: str,
    data_root: Optional[Path] = None,
    page_number: Optional[int] = None,
    extra_metadata: Optional[dict] = None,
) -> List[Document]:
    """
    构建单个标准 Document，空文本返回空列表
    Args:
        path: 源文件路径
        text: 文档文本
        file_type: 文件类型
        extra_metadata: 额外元数据字典
    Returns:
        包含单个 Document 的列表
    """
    normalized = normalize_text(text)
    if not normalized:
        return []
    metadata = build_base_metadata(
        path,
        normalized,
        file_type,
        knowledge_base_id=knowledge_base_id,
        page_number=page_number,
        data_root=data_root,
    )
    if extra_metadata:
        metadata.update(extra_metadata)
    return [Document(text=normalized, metadata=metadata)]


def _flatten_json(value: Any, prefix: str = "") -> Iterable[tuple[str, Any]]:
    """
    将嵌套 JSON 展平成 key path/value 对
    Args:
        value: 待展平的 JSON 值
        prefix: 当前路径前缀，默认为空字符串
    Returns:
        生成器，生成 key path/value 对
    """
    if isinstance(value, dict):
        # 处理字典项
        for key, item in value.items():
            next_prefix = f"{prefix}.{key}" if prefix else str(key)
            # 递归处理子项
            yield from _flatten_json(item, next_prefix)
    elif isinstance(value, list):
        for index, item in enumerate(value):
            # 处理列表项
            next_prefix = f"{prefix}[{index}]" if prefix else f"[{index}]"
            # 递归处理子项
            yield from _flatten_json(item, next_prefix)
    else:
        yield prefix, value


class _HTMLToMarkdownParser(HTMLParser):
    """轻量 HTML 到 Markdown 转换器，保留标题、列表和表格"""

    def __init__(self):
        # 将 HTML 实体转换为 Unicode 字符
        super().__init__(convert_charrefs=True)
        # 存储转换后的 Markdown 文本
        self.parts: List[str] = []
        # 存储当前标签栈
        self._stack: List[str] = []
        # 存储当前文本缓冲区
        self._buffer: List[str] = []
        # 存储当前表格
        self._table: List[List[str]] = []
        # 存储当前表格行
        self._row: List[str] = []
        # 存储当前表格单元格
        self._cell: List[str] = []

    def handle_starttag(self, tag: str, attrs: list[tuple[str, Optional[str]]]) -> None:
        """
        处理 HTML 开始标签
        Args:
            tag: 标签名称
            attrs: 标签属性列表
        """
        self._stack.append(tag)
        # 处理表格行
        if tag == "tr":
            # 新行，重置当前行
            self._row = []
        # 处理表格单元格
        elif tag in {"td", "th"}:
            # 新单元格，重置当前单元格
            self._cell = []

    def handle_endtag(self, tag: str) -> None:
        """
        处理 HTML 结束标签
        Args:
            tag: 标签名称
        """
        # 处理文本缓冲区
        text = normalize_text(" ".join(self._buffer))
        # 处理标题
        if tag in {"h1", "h2", "h3", "h4", "h5", "h6"} and text:
            # 将标题转换为 Markdown 格式
            level = int(tag[1])
            self.parts.append(f"{'#' * level} {text}")
            self._buffer = []
        # 处理段落
        elif tag == "p" and text:
            # 将段落转换为 Markdown 格式
            self.parts.append(text)
            self._buffer = []
        # 处理列表项
        elif tag == "li" and text:
            # 将列表项转换为 Markdown 格式
            self.parts.append(f"- {text}")
            self._buffer = []
        # 处理表格单元格
        elif tag in {"td", "th"}:
            # 将表格单元格内容转换为 Markdown 格式
            # 合并当前单元格内容，存储到当前行中
            self._row.append(normalize_text(" ".join(self._cell)))
            self._cell = []
        # 处理表格行
        elif tag == "tr":
            # 将表格行内容转换为 Markdown 格式
            if any(cell for cell in self._row):
                self._table.append(self._row)
            self._row = []
        # 处理表格
        elif tag == "table":
            # 将表格内容转换为 Markdown 格式
            table_text = _markdown_table(self._table)
            if table_text:
                self.parts.append(table_text)
            self._table = []
        # 处理其他标签
        if tag in self._stack:
            # 从栈中移除当前标签
            self._stack.remove(tag)

    def handle_data(self, data: str) -> None:
        """
        处理 HTML 文本数据
        Args:
            data: 待处理的文本数据
        """
        if not data or not data.strip():
            return
        current = self._stack[-1] if self._stack else ""
        # 如果当前标签是表格单元格或表头单元格，将文本添加到当前单元格
        if current in {"td", "th"}:
            self._cell.append(data.strip())
        # 否则，将文本添加到缓冲区
        else:
            self._buffer.append(data.strip())


class ParserRegistry:
    """文档解析器注册表，将文件后缀映射到解析函数"""

    def __init__(self):
        """
        初始化空解析器注册表
        """
        # 存储解析器映射
        self._parsers: Dict[str, Callable[..., List[Document]]] = {}

    @classmethod
    def default(cls) -> "ParserRegistry":
        """
        创建注册了默认文档解析器的 ParserRegistry
        Returns:
            支持 md、txt、pdf、docx 的解析器注册表
        """
        # 初始化注册表
        registry = cls()
        registry.register(".md", parse_text_like)
        registry.register(".txt", parse_text_like)
        registry.register(".csv", parse_csv)
        registry.register(".json", parse_json)
        registry.register(".html", parse_html)
        registry.register(".htm", parse_html)
        registry.register(".pdf", parse_pdf)
        registry.register(".docx", parse_docx)
        registry.register(".xlsx", parse_xlsx)
        registry.register(".pptx", parse_pptx)
        return registry

    def register(self, suffix: str, parser: Callable[..., List[Document]]) -> None:
        """
        注册指定文件后缀的解析函数
        Args:
            suffix: 文件后缀，如 .pdf
            parser: 返回 Document 列表的解析函数
        Returns:
            无返回值
        """
        self._parsers[suffix.lower()] = parser

    def load(
        self,
        path: Path,
        *,
        knowledge_base_id: str = "default",
        data_root: Optional[Path] = None,
    ) -> List[Document]:
        """
        根据文件后缀选择解析器并加载文档
        Args:
            path: 待解析文件路径
            knowledge_base_id: 资料库 ID
            data_root: 数据根目录
        Returns:
            解析得到的 Document 列表；不支持的文件返回空列表
        """
        suffix = path.suffix.lower()
        parser = self._parsers.get(suffix) # 获取文件后缀对应的解析函数
        if parser is None:
            return []
        return parser(
            path,
            knowledge_base_id=knowledge_base_id,
            data_root=data_root
        )


def parse_text_like(
    path: Path,
    *,
    knowledge_base_id: str,
    data_root: Optional[Path] = None,
) -> List[Document]:
    """
    解析 Markdown/TXT 等纯文本类文档
    Args:
        path: 文件路径
        knowledge_base_id: 资料库 ID
        data_root: 数据根目录
    Returns:
        包含全文文本和元数据的 Document 列表
    """
    file_type = path.suffix.lower().lstrip(".")
    return _document_from_text(
        path,
        read_text_file(path),
        file_type,
        knowledge_base_id=knowledge_base_id,
        data_root=data_root,
    )


def parse_csv(
    path: Path,
    *,
    knowledge_base_id: str,
    data_root: Optional[Path] = None,
) -> List[Document]:
    """
    解析 CSV 为 Markdown 表格
    Args:
        path: 文件路径
        knowledge_base_id: 资料库 ID
        data_root: 数据根目录
    Returns:
        包含表格内容的 Document 列表
    """
    content = read_text_file(path)
    if not content:
        return []
    try:
        # 读取 CSV 内容，解析为行列表
        rows = list(csv.reader(io.StringIO(content)))
    except csv.Error as exc:
        logger.warning("CSV 读取失败: %s, error=%s", path, exc)
        return []
    # 生成 Markdown 表格文本
    text = f"# {path.name}\n\n{_markdown_table(rows)}"
    # 生成 Document 列表
    return _document_from_text(path, text, "csv", knowledge_base_id=knowledge_base_id, data_root=data_root)


def parse_json(
    path: Path,
    *,
    knowledge_base_id: str,
    data_root: Optional[Path] = None,
) -> List[Document]:
    """
    解析 JSON 为可读 Markdown 键值列表
    Args:
        path: 文件路径
        knowledge_base_id: 资料库 ID
        data_root: 数据根目录
    Returns:
        包含键值对的 Document 列表
    """
    content = read_text_file(path)
    if not content:
        return []
    try:
        # 读取 JSON 内容，解析为 Python 对象
        payload = json.loads(content)
    except json.JSONDecodeError as exc:
        logger.warning("JSON 读取失败: %s, error=%s", path, exc)
        return []
    lines = [f"# {path.name}"]
    for key, value in _flatten_json(payload):
        lines.append(f"- {key}: {value}")
    return _document_from_text(path, "\n".join(lines), "json", knowledge_base_id=knowledge_base_id, data_root=data_root)


def parse_html(
    path: Path,
    *,
    knowledge_base_id: str,
    data_root: Optional[Path] = None,
) -> List[Document]:
    """
    解析 HTML/HTM，保留标题、段落、列表和表格
    Args:
        path: 文件路径
        knowledge_base_id: 资料库 ID
        data_root: 数据根目录
    Returns:
        包含解析结果的 Document 列表
    """
    content = read_text_file(path)
    if not content:
        return []
    try:
        from bs4 import BeautifulSoup
    except Exception:
        soup = None
    else:
        # 使用 BeautifulSoup 解析 HTML 内容，生成可操作的文档树对象
        soup = BeautifulSoup(content, "html.parser")

    if soup is not None:
        parts: List[str] = []
        # 遍历文档树，提取标题、段落、列表和表格
        for element in soup.find_all(["h1", "h2", "h3", "h4", "h5", "h6", "p", "li", "table"]):
            name = element.name.lower()
            # 处理标题元素
            if name.startswith("h"):
                parts.append(f"{'#' * int(name[1])} {normalize_text(element.get_text(' ', strip=True))}")
            # 处理列表元素
            elif name == "li":
                parts.append(f"- {normalize_text(element.get_text(' ', strip=True))}")
            # 处理段落元素
            elif name == "p":
                parts.append(normalize_text(element.get_text(" ", strip=True)))
            # 处理表格元素
            elif name == "table":
                rows = [
                    [cell.get_text(" ", strip=True) for cell in row.find_all(["th", "td"])]
                    for row in element.find_all("tr")
                ]
                table_text = _markdown_table(rows)
                if table_text:
                    parts.append(table_text)
        return _document_from_text(path, "\n\n".join(parts), path.suffix.lower().lstrip("."), knowledge_base_id=knowledge_base_id, data_root=data_root)

    # 处理其他情况，如解析失败或未安装 BeautifulSoup 库
    parser = _HTMLToMarkdownParser()
    # 解析 HTML 内容，生成 Markdown 文本
    parser.feed(content)
    return _document_from_text(path, "\n\n".join(parser.parts), path.suffix.lower().lstrip("."), knowledge_base_id=knowledge_base_id, data_root=data_root)


def parse_pdf(
    path: Path,
    *,
    knowledge_base_id: str,
    data_root: Optional[Path] = None,
) -> List[Document]:
    """
    按页解析 PDF 文档，并为每页生成独立 Document
    Args:
        path: PDF 文件路径
        knowledge_base_id: 资料库 ID
        data_root: 数据根目录
    Returns:
        每个非空页面对应的 Document 列表
    """
    pages = _extract_pdf_pages_with_pymupdf(path)
    if not pages:
        pages = _extract_pdf_pages_with_pypdf(path)
    return [
        Document(
            text=text,
            metadata=build_base_metadata(
                path,
                text,
                "pdf",
                knowledge_base_id=knowledge_base_id,
                page_number=page_number,
                data_root=data_root,
            ),
        )
        for page_number, text in pages
    ]


def _extract_pdf_pages_with_pymupdf(path: Path) -> List[tuple[int, str]]:
    """
    优先使用 PyMuPDF 提取 PDF 文本，表格/版面保真度通常优于 pypdf
    Args:
        path: PDF 文件路径
    Returns:
        每个非空页面对应的文本元组列表
    """
    try:
        import fitz
    except Exception:
        return []
    try:
        # 基于 PyMuPDF（fitz） 库打开本地 pdf 文件，创建 pdf 文档对象
        pdf = fitz.open(str(path))
    except Exception as exc:
        logger.warning("PDF 读取失败: %s, error=%s", path, exc)
        return []
    pages: List[tuple[int, str]] = []
    try:
        for index, page in enumerate(pdf, 1):
            text = normalize_text(page.get_text("text") or "")
            if text:
                pages.append((index, text))
    finally:
        pdf.close()
    return pages


def _extract_pdf_pages_with_pypdf(path: Path) -> List[tuple[int, str]]:
    """
    使用 pypdf 作为 PDF 文本提取兜底
    Args:
        path: PDF 文件路径
    Returns:
        每个非空页面对应的文本元组列表
    """
    try:
        # 基于 pypdf 库打开本地 pdf 文件，创建 pdf 文档对象
        reader = PdfReader(str(path))
    except Exception as exc:
        logger.warning("PDF 读取失败: %s, error=%s", path, exc)
        return []
    pages = []
    for page_number, page in enumerate(reader.pages, 1):
        text = normalize_text(page.extract_text() or "")
        if text:
            pages.append((page_number, text))
    return pages


def parse_docx(
    path: Path,
    *,
    knowledge_base_id: str,
    data_root: Optional[Path] = None,
) -> List[Document]:
    """
    解析 DOCX 文档段落文本
    Args:
        path: DOCX 文件路径
        knowledge_base_id: 资料库 ID
        data_root: 数据根目录
    Returns:
        包含全文文本和元数据的 Document 列表
    """
    try:
        from docx import Document as DocxDocument
        from docx.oxml.table import CT_Tbl
        from docx.oxml.text.paragraph import CT_P
        from docx.table import Table
        from docx.text.paragraph import Paragraph
    except Exception as exc:
        logger.warning("DOCX 解析依赖不可用: %s", exc)
        return []
    try:
        docx = DocxDocument(str(path))
    except Exception as exc:
        logger.warning("DOCX 读取失败: %s, error=%s", path, exc)
        return []
    parts: List[str] = []
    # 按文档书写顺序遍历正文内每一个子元素
    for child in docx.element.body.iterchildren():
        # 如果子元素是段落
        if isinstance(child, CT_P):
            # 将底层xml节点包装为段落对象，提取文本
            text = Paragraph(child, docx).text
            if normalize_text(text):
                parts.append(text)
        # 如果子元素是表格
        elif isinstance(child, CT_Tbl):
            # 将底层xml节点包装为表格对象，提取文本
            table = Table(child, docx)
            # 提取表格所有单元格文本，过滤掉空单元格
            rows = [[cell.text for cell in row.cells] for row in table.rows]
            # 转换为 Markdown 表格格式
            table_text = _markdown_table(rows)
            if table_text:
                parts.append(table_text)
    return _document_from_text(path, "\n\n".join(parts), "docx", knowledge_base_id=knowledge_base_id, data_root=data_root)


def parse_xlsx(
    path: Path,
    *,
    knowledge_base_id: str,
    data_root: Optional[Path] = None,
) -> List[Document]:
    """
    解析 XLSX 工作簿，每个工作表转为 Markdown 表格
    Args:
        path: XLSX 文件路径
        knowledge_base_id: 资料库 ID
        data_root: 数据根目录
    Returns:
        包含工作表文本和元数据的 Document 列表
    """
    try:
        import openpyxl
    except Exception as exc:
        logger.warning("XLSX 解析依赖不可用: %s", exc)
        return []
    try:
        # 加载工作簿，只读模式，只解析数据，不解析公式
        workbook = openpyxl.load_workbook(str(path), data_only=True, read_only=True)
    except Exception as exc:
        logger.warning("XLSX 读取失败: %s, error=%s", path, exc)
        return []
    parts: List[str] = [f"# {path.name}"]
    try:
        # 遍历工作簿中的每个工作表
        for sheet in workbook.worksheets:
            # 提取工作表所有单元格文本，过滤掉空单元格
            rows = [list(row) for row in sheet.iter_rows(values_only=True)]
            # 转换为 Markdown 表格格式
            table_text = _markdown_table(rows)
            if table_text:
                parts.append(f"## {sheet.title}\n\n{table_text}")
    finally:
        workbook.close()
    return _document_from_text(path, "\n\n".join(parts), "xlsx", knowledge_base_id=knowledge_base_id, data_root=data_root)


def parse_pptx(
    path: Path,
    *,
    knowledge_base_id: str,
    data_root: Optional[Path] = None,
) -> List[Document]:
    """
    解析 PPTX 幻灯片文本和表格为 Markdown 表格
    Args:
        path: PPTX 文件路径
        knowledge_base_id: 资料库 ID
        data_root: 数据根目录
    Returns:
        包含幻灯片文本和元数据的 Document 列表
    """
    try:
        from pptx import Presentation
    except Exception as exc:
        logger.warning("PPTX 解析依赖不可用: %s", exc)
        return []
    try:
        # 加载演示文稿
        presentation = Presentation(str(path))
    except Exception as exc:
        logger.warning("PPTX 读取失败: %s, error=%s", path, exc)
        return []
    parts: List[str] = [f"# {path.name}"]
    # 遍历演示文稿中的每个幻灯片
    for index, slide in enumerate(presentation.slides, 1):
        slide_parts = [f"## Slide {index}"]
        for shape in slide.shapes:
            # 如果子元素是表格
            if getattr(shape, "has_table", False):
                # 提取表格所有单元格文本，过滤掉空单元格
                rows = [[cell.text for cell in row.cells] for row in shape.table.rows]
                # 转换为 Markdown 表格格式
                table_text = _markdown_table(rows)
                if table_text:
                    slide_parts.append(table_text)
            # 如果子元素是文本框
            elif getattr(shape, "has_text_frame", False):
                text = normalize_text(shape.text)
                if text:
                    slide_parts.append(text)
        if len(slide_parts) > 1:
            parts.append("\n\n".join(slide_parts))
    return _document_from_text(path, "\n\n".join(parts), "pptx", knowledge_base_id=knowledge_base_id, data_root=data_root)


def load_documents(
    data_path: str | Path,
    *,
    knowledge_base_id: str = "default",
    document_registry: Optional[object] = None,
    parser_registry: Optional[ParserRegistry] = None,
) -> List[Document]:
    """
    从数据目录加载支持格式文档，或按文档注册表加载未删除文档
    Args:
        data_path: 数据根目录
        knowledge_base_id: 默认资料库 ID
        document_registry: 可选文档生命周期注册表
        parser_registry: 可选解析器注册表
    Returns:
        解析得到的 LlamaIndex Document 列表
    """
    root = Path(data_path)
    # 初始化解析器注册表
    registry = parser_registry or ParserRegistry.default()
    if not root.exists():
        logger.warning("文档目录不存在: %s", root)
        return []
    docs: List[Document] = []
    if document_registry is not None:
        records = []
        for record in document_registry.records.values():
            path = Path(record.source_path)
            if record.knowledge_base_id != knowledge_base_id:
                continue
            # 跳过已删除或文件不存在的注册文档
            if record.status == "deleted" or not path.exists():
                continue
            # 跳过可信目录外的注册文档
            if not is_path_within(path, root):
                logger.warning("跳过可信目录外的注册文档: %s", path)
                continue
            records.append(record)
        # 按上传时间排序后加载注册文档，确保先上传的文档先被加载和解析
        for record in sorted(records, key=lambda item: item.upload_time):
            path = Path(record.source_path)
            # 仅加载支持格式的注册文档
            if path.is_file() and path.suffix.lower() in SUPPORTED_SUFFIXES:
                # 加载注册文档
                parsed_docs = registry.load(
                    path,
                    knowledge_base_id=record.knowledge_base_id,
                    data_root=root,
                )
                for doc in parsed_docs:
                    # 将注册表中的文档元数据更新到解析得到的 Document 中，确保与注册表一致
                    doc.metadata.update(
                        {
                            "knowledge_base_id": record.knowledge_base_id,
                            "document_id": record.document_id,
                            "filename": record.filename,
                            "file_type": record.file_type,
                        }
                    )
                docs.extend(parsed_docs)
        return docs
    # 无注册表时按文件系统加载，可能包含已删除或未注册的文档
    for path in sorted(root.rglob("*")): # 递归遍历数据根目录
        if path.is_file() and path.suffix.lower() in SUPPORTED_SUFFIXES:
            docs.extend(
                registry.load(
                    path,
                    knowledge_base_id=knowledge_base_id,
                    data_root=root,
                )
            )
    return docs


@dataclass
class DocumentRecord:
    """文档生命周期注册记录"""

    document_id: str # 文档 ID
    source_path: str # 源文件路径
    filename: str # 文件名
    file_type: str # 文件类型，如 pdf、md
    content_hash: str # 文件内容哈希值
    knowledge_base_id: str # 资料库 ID
    status: str = "uploaded" # 文档状态，如 uploaded、parsed、deleted
    chunk_count: int = 0 # 该文档生成的叶子分块数量
    upload_time: str = "" # 文档上传时间 ISO 字符串
    last_error: Optional[str] = None # 最新解析错误信息

    def to_dict(self) -> dict:
        """
        将文档记录转换为可 JSON 序列化字典
        Returns:
            文档生命周期记录字典
        """
        return asdict(self)
